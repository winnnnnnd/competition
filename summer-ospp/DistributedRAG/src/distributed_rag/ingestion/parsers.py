from __future__ import annotations

import html
import io
import os
import re
import subprocess
import tempfile
from dataclasses import asdict
from pathlib import Path
from typing import Any, Dict, Iterator, List, Sequence

from ..config import ParsingConfig, StorageConfig
from ..domain.identifiers import deterministic_object_key, make_element_id
from ..domain.models import Document, DocumentElement, ElementType, FileType, SourceLocator
from ..infrastructure.object_storage import ObjectStorage


def _batched(values: Sequence[DocumentElement], size: int) -> Iterator[List[Dict[str, Any]]]:
    for start in range(0, len(values), size):
        yield [value.as_dict() for value in values[start:start + size]]


def _element(document: Document, ordinal: int, element_type: ElementType, text: str, locator: SourceLocator, metadata: Dict[str, Any] | None = None, parent_id: str | None = None) -> DocumentElement:
    return DocumentElement(
        element_id=make_element_id(document.document_version, element_type.value, locator, ordinal),
        document_id=document.document_id,
        document_version=document.document_version,
        element_type=element_type.value,
        text=text.strip(),
        source_locator=locator,
        parent_id=parent_id,
        metadata={"file_name": document.file_name, **document.metadata, **(metadata or {})},
    )


def parse_document_batches(document_value: Dict[str, Any], storage_value: Dict[str, Any], parsing_value: Dict[str, Any]) -> Iterator[List[Dict[str, Any]]]:
    """Parse a MinIO object and yield bounded element batches from a Ray generator task."""
    document = Document(**document_value)
    storage = ObjectStorage(StorageConfig(**storage_value))
    config = ParsingConfig(**parsing_value)
    file_type = FileType(document.file_type)

    if file_type == FileType.PDF:
        yield from _parse_pdf_batches(document, storage, config)
        return
    if file_type == FileType.IMAGE:
        elements = [_element(document, 0, ElementType.OCR_TEXT, "", SourceLocator(page_number=1), {"requires_ocr": True, "ocr_uri": document.source_uri})]
    elif file_type == FileType.AUDIO:
        elements = [_element(document, 0, ElementType.ASR_SEGMENT, "", SourceLocator(), {"requires_asr": True, "asr_uri": document.source_uri})]
    else:
        content = storage.get_bytes(document.source_uri)
        if file_type == FileType.WORD:
            elements = _parse_word(content, document)
        elif file_type == FileType.POWERPOINT:
            elements = _parse_powerpoint(content, document)
        elif file_type == FileType.EXCEL:
            elements = _parse_excel(content, document)
        elif file_type == FileType.MARKDOWN:
            elements = _parse_markdown(content, document)
        elif file_type in {FileType.TEXT, FileType.WEB}:
            elements = _parse_text(content, document)
        else:
            raise ValueError(f"Unsupported file type: {document.file_name}")

    yield from _batched(elements, config.element_batch_size)


def _parse_pdf_batches(document: Document, storage: ObjectStorage, config: ParsingConfig) -> Iterator[List[Dict[str, Any]]]:
    from PyPDF2 import PdfReader

    with tempfile.NamedTemporaryFile(prefix="rag-pdf-", suffix=".pdf") as local_file:
        storage.download_to_file(document.source_uri, local_file)
        local_file.flush()
        reader = PdfReader(local_file.name)
        batch: List[DocumentElement] = []
        for page_number, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text() or ""
            if page_text.strip():
                locator = SourceLocator(page_number=page_number, text_start=0, text_end=len(page_text))
                batch.append(_element(document, page_number - 1, ElementType.TEXT, page_text, locator))
            elif config.enable_ocr_fallback:
                image_uri = _render_pdf_page(local_file.name, page_number, document, storage)
                locator = SourceLocator(page_number=page_number)
                batch.append(_element(document, page_number - 1, ElementType.OCR_TEXT, "", locator, {"requires_ocr": True, "ocr_uri": image_uri}))
            if len(batch) >= config.element_batch_size:
                yield [element.as_dict() for element in batch]
                batch = []
        if batch:
            yield [element.as_dict() for element in batch]


def _render_pdf_page(local_path: str, page_number: int, document: Document, storage: ObjectStorage) -> str:
    import fitz

    pdf = fitz.open(local_path)
    try:
        page = pdf.load_page(page_number - 1)
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        image = pixmap.tobytes("png")
    finally:
        pdf.close()
    key = deterministic_object_key(document.document_id, document.document_version, "ocr-pages", f"page-{page_number:06d}.png")
    return storage.put_bytes(key, image, "image/png")


def _convert_legacy_office(content: bytes, suffix: str, target_suffix: str) -> bytes:
    with tempfile.TemporaryDirectory(prefix="rag-office-") as temp_dir:
        source = Path(temp_dir) / f"source{suffix}"
        source.write_bytes(content)
        completed = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", target_suffix.lstrip("."), "--outdir", temp_dir, str(source)],
            check=False,
            capture_output=True,
            timeout=120,
        )
        target = source.with_suffix(target_suffix)
        if completed.returncode != 0 or not target.exists():
            raise ValueError(f"LibreOffice conversion failed: {completed.stderr.decode('utf-8', errors='ignore')}")
        return target.read_bytes()


def _parse_word(content: bytes, document: Document) -> List[DocumentElement]:
    import docx

    if document.file_name.lower().endswith(".doc"):
        content = _convert_legacy_office(content, ".doc", ".docx")
    word = docx.Document(io.BytesIO(content))
    elements: List[DocumentElement] = []
    headings: List[str] = []
    for paragraph_number, paragraph in enumerate(word.paragraphs, start=1):
        text = paragraph.text.strip()
        if not text:
            continue
        style = paragraph.style.name if paragraph.style else ""
        is_heading = style.lower().startswith("heading")
        if is_heading:
            match = re.search(r"(\d+)", style)
            level = int(match.group(1)) if match else 1
            headings = headings[:level - 1] + [text]
        locator = SourceLocator(paragraph_number=paragraph_number, heading_path=list(headings))
        elements.append(_element(document, paragraph_number, ElementType.TITLE if is_heading else ElementType.TEXT, text, locator, {"style": style}))
    return elements


def _parse_powerpoint(content: bytes, document: Document) -> List[DocumentElement]:
    import pptx

    if document.file_name.lower().endswith(".ppt"):
        content = _convert_legacy_office(content, ".ppt", ".pptx")
    presentation = pptx.Presentation(io.BytesIO(content))
    elements: List[DocumentElement] = []
    ordinal = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape_number, shape in enumerate(slide.shapes, start=1):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
            if not text.strip():
                continue
            bbox = [float(shape.left), float(shape.top), float(shape.left + shape.width), float(shape.top + shape.height)]
            locator = SourceLocator(slide_number=slide_number, bbox=bbox)
            elements.append(_element(document, ordinal, ElementType.TEXT, text, locator, {"shape_number": shape_number, "shape_name": shape.name}))
            ordinal += 1
    return elements


def _parse_excel(content: bytes, document: Document) -> List[DocumentElement]:
    import openpyxl
    from openpyxl.utils import get_column_letter

    if document.file_name.lower().endswith(".csv"):
        return _parse_csv(content, document)
    workbook = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    elements: List[DocumentElement] = []
    ordinal = 0
    try:
        for worksheet in workbook.worksheets:
            rows: List[List[str]] = []
            start_row = 1
            for row_number, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                values = ["" if value is None else str(value) for value in row]
                if not any(values):
                    continue
                if not rows:
                    start_row = row_number
                rows.append(values)
                if len(rows) >= 100:
                    ordinal = _append_sheet_block(elements, document, worksheet.title, rows, start_row, row_number, ordinal)
                    rows = []
            if rows:
                ordinal = _append_sheet_block(elements, document, worksheet.title, rows, start_row, start_row + len(rows) - 1, ordinal)
    finally:
        workbook.close()
    return elements


def _append_sheet_block(elements: List[DocumentElement], document: Document, sheet_name: str, rows: List[List[str]], start_row: int, end_row: int, ordinal: int) -> int:
    from openpyxl.utils import get_column_letter

    width = max((len(row) for row in rows), default=1)
    cell_range = f"A{start_row}:{get_column_letter(width)}{end_row}"
    text = "\n".join("\t".join(row) for row in rows)
    locator = SourceLocator(sheet_name=sheet_name, cell_range=cell_range)
    elements.append(_element(document, ordinal, ElementType.TABLE, text, locator))
    return ordinal + 1


def _parse_csv(content: bytes, document: Document) -> List[DocumentElement]:
    import csv

    decoded = content.decode("utf-8-sig", errors="replace")
    rows = list(csv.reader(io.StringIO(decoded)))
    elements: List[DocumentElement] = []
    for start in range(0, len(rows), 100):
        block = rows[start:start + 100]
        _append_sheet_block(elements, document, "CSV", block, start + 1, start + len(block), len(elements))
    return elements


def _parse_markdown(content: bytes, document: Document) -> List[DocumentElement]:
    text = content.decode("utf-8", errors="replace")
    elements: List[DocumentElement] = []
    headings: List[str] = []
    buffer: List[str] = []
    start_offset = 0

    def flush(end_offset: int) -> None:
        nonlocal buffer, start_offset
        value = "\n".join(buffer).strip()
        if value:
            locator = SourceLocator(heading_path=list(headings), text_start=start_offset, text_end=end_offset)
            elements.append(_element(document, len(elements), ElementType.TEXT, value, locator))
        buffer = []

    offset = 0
    for line in text.splitlines(keepends=True):
        match = re.match(r"^(#{1,6})\s+(.+?)\s*$", line)
        if match:
            flush(offset)
            level = len(match.group(1))
            title = match.group(2)
            headings = headings[:level - 1] + [title]
            locator = SourceLocator(heading_path=list(headings), text_start=offset, text_end=offset + len(line))
            elements.append(_element(document, len(elements), ElementType.TITLE, title, locator))
            start_offset = offset + len(line)
        else:
            if not buffer:
                start_offset = offset
            buffer.append(line.rstrip("\n"))
        offset += len(line)
    flush(len(text))
    return elements


def _parse_text(content: bytes, document: Document) -> List[DocumentElement]:
    text = content.decode("utf-8", errors="replace")
    if document.file_name.lower().endswith((".html", ".htm")):
        text = re.sub(r"<script.*?</script>|<style.*?</style>", " ", text, flags=re.I | re.S)
        text = html.unescape(re.sub(r"<[^>]+>", " ", text))
        text = re.sub(r"\s+", " ", text)
    locator = SourceLocator(text_start=0, text_end=len(text))
    return [_element(document, 0, ElementType.TEXT, text, locator)] if text.strip() else []
